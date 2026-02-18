import os
import pypdf
from docx import Document
from pptx import Presentation
from flask import current_app
import logging

# Configure Logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

class DocumentEditor:
    """
    Advanced Document Processor.
    Reads: PDF, DOCX, PPTX.
    Writes: DOCX, TXT.
    """

    def __init__(self):
        pass

    # ==========================
    # 📖 READING (Extraction)
    # ==========================

    def read_document(self, filepath, file_type):
        """
        Smart router to extract text based on file type.
        """
        if not os.path.exists(filepath):
            return "Error: File not found."

        try:
            if file_type == 'pdf':
                return self._read_pdf(filepath)
            elif file_type == 'docx':
                return self._read_docx(filepath)
            elif file_type == 'pptx':
                return self._read_pptx(filepath)
            elif file_type == 'txt' or file_type == 'md':
                with open(filepath, 'r', encoding='utf-8') as f:
                    return f.read()
            else:
                return "Error: Unsupported document format for reading."
        except Exception as e:
            logger.error(f"Read error: {e}")
            return f"Could not read document: {str(e)}"

    def _read_pdf(self, path):
        """Extracts text from PDF."""
        text = ""
        reader = pypdf.PdfReader(path)
        for page in reader.pages:
            text += page.extract_text() + "\n"
        return text

    def _read_docx(self, path):
        """Extracts text from Word Doc."""
        doc = Document(path)
        return "\n".join([para.text for para in doc.paragraphs])

    def _read_pptx(self, path):
        """Extracts text from PowerPoint slides."""
        prs = Presentation(path)
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text

    # ==========================
    # ✍️ WRITING (Generation)
    # ==========================

    def create_document(self, content, filename="generated_doc.docx"):
        """
        Creates a Word Document from AI text and saves it.
        """
        try:
            doc = Document()
            doc.add_heading('AI Generated Report', 0)
            
            # Split by newlines and add paragraphs
            for line in content.split('\n'):
                if line.strip():
                    doc.add_paragraph(line)
            
            # Save
            folder = os.path.join(current_app.config['UPLOAD_FOLDER'], 'docs')
            if not os.path.exists(folder):
                os.makedirs(folder)
                
            path = os.path.join(folder, filename)
            doc.save(path)
            
            return path
        except Exception as e:
            logger.error(f"Write error: {e}")
            return None

# Global Instance
doc_editor = DocumentEditor()